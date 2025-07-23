import io, tempfile, csv, json
from typing import Optional
from pypdf import PdfReader
import docx2txt
import pandas as pd
from pptx import Presentation

class PasswordRequired(Exception):
    pass

def load_pdf(content: bytes, password: Optional[str]=None) -> str:
    reader = PdfReader(io.BytesIO(content))
    if reader.is_encrypted:
        if not password:
            raise PasswordRequired("PDF is encrypted")
        try:
            reader.decrypt(password)
        except Exception:
            raise PasswordRequired("Wrong password")
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)

def load_docx(content: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".docx") as f:
        f.write(content); f.flush()
        return docx2txt.process(f.name)

def load_pptx(content: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".pptx") as f:
        f.write(content); f.flush()
        prs = Presentation(f.name)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def load_excel(content: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".xlsx") as f:
        f.write(content); f.flush()
        df = pd.read_excel(f.name, sheet_name=None)
    parts = []
    for name, sheet in df.items():
        parts.append(f"# sheet: {name}\n" + sheet.to_csv(index=False, sep="\t"))
    return "\n".join(parts)

def load_csv(content: bytes) -> str:
    return content.decode("utf-8", errors="ignore")

def load_json(content: bytes) -> str:
    try:
        obj = json.loads(content.decode("utf-8", errors="ignore"))
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        return content.decode("utf-8", errors="ignore")

def load_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="ignore")

EXT_LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_excel,
    ".csv": load_csv,
    ".json": load_json,
    ".txt": load_txt,
    ".md": load_txt,
    ".html": load_txt,
}
